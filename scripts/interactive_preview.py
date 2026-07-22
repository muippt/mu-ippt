#!/usr/bin/env python3
"""
Interactive PPTX preview: generate a self-contained HTML page that renders
every slide as an image with clickable shape overlays.

Inspired by OfficeCLI's `watch` + `get selected` interaction model: instead of
reading raw XML, a human (or an agent driving a headless browser) can click on
a rendered slide and instantly see which shape occupies that region, along
with its address (`/slide[N]/shape[M]`), name, type, text content and
position/size.

Rendering strategy (best available wins):
    1. Reuse existing JPEG slide images produced by visual_verify.py
       (`<pptx_dir>/visual_review/slide-*.jpg`) if present.
    2. Otherwise call LibreOffice (`soffice`) + `pdftoppm` to rasterize the
       deck (same pipeline as visual_verify.py).
    3. Otherwise fall back to a pure python-pptx + Pillow rendering that
       draws simple placeholder boxes for every shape (position + label),
       so the tool always produces *something* usable, even with zero
       external dependencies beyond python-pptx/Pillow.

Usage:
    python3 interactive_preview.py <input.pptx> [--output <html_path>] [--port 0]

    --output PATH   Where to write the HTML file.
                    Default: <pptx_dir>/<pptx_stem>_preview.html
    --port N        If given, start a simple HTTP server on port N (use 0 for
                    an OS-assigned free port) serving the directory that
                    contains the generated HTML/images, and print the URL.
                    Press Ctrl+C to stop the server.
    --dpi N         Rasterization DPI when falling back to LibreOffice
                    (default: 150).

Requires: python-pptx, Pillow. LibreOffice (soffice) + poppler (pdftoppm) are
optional but recommended for accurate rendering.
"""

from __future__ import annotations

import argparse
import base64
import html
import json
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

try:
    from PIL import Image, ImageDraw, ImageFont
    _HAS_PIL = True
except ImportError:
    _HAS_PIL = False

__version__ = "1.2.0"

EMU_PER_INCH = 914400
PLACEHOLDER_SLIDE_WIDTH_PX = 960
PLACEHOLDER_SLIDE_HEIGHT_PX = 540


# --------------------------------------------------------------------------
# Shape introspection
# --------------------------------------------------------------------------

def _shape_type_name(shape: Any) -> str:
    try:
        return str(shape.shape_type).split(".")[-1].split(" ")[0]
    except Exception:
        return shape.__class__.__name__


def _safe_text(shape: Any) -> str:
    try:
        if getattr(shape, "has_text_frame", False):
            return shape.text_frame.text
    except Exception:
        pass
    try:
        if getattr(shape, "has_table", False):
            cells = []
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        cells.append(cell.text.strip())
            return " | ".join(cells)
        if getattr(shape, "has_chart", False):
            return f"[chart: {shape.chart.chart_type}]"
    except Exception:
        pass
    return ""


def collect_shapes(
    shape: Any,
    slide_idx: int,
    path_prefix: str,
    parent_left: int,
    parent_top: int,
    depth: int = 0,
) -> List[Dict[str, Any]]:
    """Recursively collect shape metadata with absolute EMU positions.

    Path format mirrors OfficeCLI-style addressing: /slide[N]/shape[M] with
    nested groups expressed as /slide[N]/shape[M]/shape[K].
    """
    results: List[Dict[str, Any]] = []

    left = getattr(shape, "left", None) or 0
    top = getattr(shape, "top", None) or 0
    width = getattr(shape, "width", None) or 0
    height = getattr(shape, "height", None) or 0

    abs_left = parent_left + left
    abs_top = parent_top + top

    info = {
        "path": path_prefix,
        "slide_index": slide_idx,
        "name": getattr(shape, "name", "") or "",
        "shape_type": _shape_type_name(shape),
        "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
        "text": _safe_text(shape),
        "left_emu": abs_left,
        "top_emu": abs_top,
        "width_emu": width,
        "height_emu": height,
        "left_in": round(abs_left / EMU_PER_INCH, 3),
        "top_in": round(abs_top / EMU_PER_INCH, 3),
        "width_in": round(width / EMU_PER_INCH, 3),
        "height_in": round(height / EMU_PER_INCH, 3),
        "depth": depth,
    }
    results.append(info)

    # Recurse into groups
    if hasattr(shape, "shapes"):
        for idx, child in enumerate(shape.shapes, start=1):
            child_path = f"{path_prefix}/shape[{idx}]"
            results.extend(
                collect_shapes(child, slide_idx, child_path, abs_left, abs_top, depth + 1)
            )

    return results


def build_slide_index(prs: "Presentation") -> Tuple[int, int, List[List[Dict[str, Any]]]]:
    """Return (slide_width_emu, slide_height_emu, per-slide shape metadata list)."""
    slide_w = prs.slide_width or (13.333 * EMU_PER_INCH)
    slide_h = prs.slide_height or (7.5 * EMU_PER_INCH)

    all_slides: List[List[Dict[str, Any]]] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        shapes: List[Dict[str, Any]] = []
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            path = f"/slide[{slide_idx}]/shape[{shape_idx}]"
            shapes.extend(collect_shapes(shape, slide_idx, path, 0, 0))
        all_slides.append(shapes)

    return slide_w, slide_h, all_slides


# --------------------------------------------------------------------------
# Rendering strategies
# --------------------------------------------------------------------------

def find_existing_review_images(pptx_path: Path) -> List[Path]:
    """Look for JPEGs already produced by visual_verify.py."""
    review_dir = pptx_path.parent / "visual_review"
    if not review_dir.is_dir():
        return []
    images = sorted(review_dir.glob(f"{pptx_path.stem}-*.jpg")) or sorted(review_dir.glob("slide-*.jpg"))
    return images


def render_via_libreoffice(pptx_path: Path, out_dir: Path, dpi: int = 150) -> List[Path]:
    """Rasterize via soffice + pdftoppm, mirroring visual_verify.py's pipeline."""
    if shutil.which("soffice") is None or shutil.which("pdftoppm") is None:
        return []

    out_dir.mkdir(parents=True, exist_ok=True)
    pdf_path = out_dir / f"{pptx_path.stem}.pdf"

    with tempfile.TemporaryDirectory() as tmpdir:
        cmd = [
            "soffice", "--headless", "--convert-to", "pdf",
            "--outdir", tmpdir, str(pptx_path),
        ]
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
        except (subprocess.TimeoutExpired, OSError):
            return []
        if result.returncode != 0:
            return []

        tmp_pdf = Path(tmpdir) / f"{pptx_path.stem}.pdf"
        if not tmp_pdf.exists():
            return []
        shutil.move(str(tmp_pdf), str(pdf_path))

    prefix = out_dir / "preview-slide"
    cmd = ["pdftoppm", "-jpeg", "-r", str(dpi), str(pdf_path), str(prefix)]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
    except (subprocess.TimeoutExpired, OSError):
        return []
    if result.returncode != 0:
        return []

    pdf_path.unlink(missing_ok=True)
    return sorted(out_dir.glob("preview-slide-*.jpg"))


def render_placeholders(
    slide_w_emu: int,
    slide_h_emu: int,
    slides_shapes: List[List[Dict[str, Any]]],
    out_dir: Path,
) -> List[Path]:
    """Draw simple placeholder boxes for every shape using Pillow only."""
    out_dir.mkdir(parents=True, exist_ok=True)
    images: List[Path] = []

    if not _HAS_PIL:
        return images

    aspect = slide_h_emu / slide_w_emu if slide_w_emu else 0.5625
    width_px = PLACEHOLDER_SLIDE_WIDTH_PX
    height_px = max(1, int(width_px * aspect))

    try:
        font = ImageFont.load_default()
    except Exception:
        font = None

    palette = ["#4C6EF5", "#F76707", "#12B886", "#E64980", "#7048E8", "#1098AD"]

    for slide_idx, shapes in enumerate(slides_shapes, start=1):
        img = Image.new("RGB", (width_px, height_px), "#FFFFFF")
        draw = ImageDraw.Draw(img)
        draw.rectangle([0, 0, width_px - 1, height_px - 1], outline="#CCCCCC", width=2)

        for i, shape in enumerate(shapes):
            if shape["depth"] > 0:
                continue  # keep placeholder simple: top-level shapes only
            x0 = int((shape["left_emu"] / slide_w_emu) * width_px) if slide_w_emu else 0
            y0 = int((shape["top_emu"] / slide_h_emu) * height_px) if slide_h_emu else 0
            x1 = int(((shape["left_emu"] + shape["width_emu"]) / slide_w_emu) * width_px) if slide_w_emu else width_px
            y1 = int(((shape["top_emu"] + shape["height_emu"]) / slide_h_emu) * height_px) if slide_h_emu else height_px
            color = palette[i % len(palette)]
            draw.rectangle([x0, y0, max(x0 + 1, x1), max(y0 + 1, y1)], outline=color, width=2)
            label = shape["name"] or shape["shape_type"]
            draw.text((x0 + 3, y0 + 3), label[:24], fill=color, font=font)

        draw.text((6, height_px - 16), f"Slide {slide_idx} (placeholder render)", fill="#999999", font=font)

        out_path = out_dir / f"placeholder-slide-{slide_idx}.jpg"
        img.save(out_path, "JPEG", quality=88)
        images.append(out_path)

    return images


def image_to_data_uri(path: Path) -> str:
    data = path.read_bytes()
    b64 = base64.b64encode(data).decode("ascii")
    ext = path.suffix.lower().lstrip(".")
    mime = "jpeg" if ext in ("jpg", "jpeg") else ext
    return f"data:image/{mime};base64,{b64}"


# --------------------------------------------------------------------------
# HTML generation
# --------------------------------------------------------------------------

HTML_TEMPLATE = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>mu-ippt Interactive Preview — {title}</title>
<style>
  :root {{
    --bg: #0f1115;
    --panel: #171a21;
    --panel-2: #1e222b;
    --border: #2a2f3a;
    --text: #e6e8ec;
    --text-dim: #9aa1ac;
    --accent: #5b8cff;
    --accent-2: #ffb454;
    --radius: 10px;
  }}
  * {{ box-sizing: border-box; }}
  body {{
    margin: 0;
    font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, Helvetica, Arial, sans-serif;
    background: var(--bg);
    color: var(--text);
    height: 100vh;
    display: flex;
    flex-direction: column;
  }}
  header {{
    display: flex;
    align-items: center;
    gap: 16px;
    padding: 10px 18px;
    background: var(--panel);
    border-bottom: 1px solid var(--border);
    flex-shrink: 0;
  }}
  header h1 {{
    font-size: 15px;
    font-weight: 600;
    margin: 0;
    color: var(--text);
    white-space: nowrap;
    overflow: hidden;
    text-overflow: ellipsis;
  }}
  header .meta {{
    color: var(--text-dim);
    font-size: 12px;
  }}
  header .nav {{
    margin-left: auto;
    display: flex;
    align-items: center;
    gap: 8px;
  }}
  button {{
    background: var(--panel-2);
    color: var(--text);
    border: 1px solid var(--border);
    border-radius: 6px;
    padding: 6px 12px;
    font-size: 13px;
    cursor: pointer;
  }}
  button:hover {{ border-color: var(--accent); }}
  button:disabled {{ opacity: 0.4; cursor: not-allowed; }}
  #page-indicator {{
    font-variant-numeric: tabular-nums;
    color: var(--text-dim);
    font-size: 13px;
  }}
  main {{
    flex: 1;
    display: flex;
    min-height: 0;
  }}
  #sidebar {{
    width: 220px;
    flex-shrink: 0;
    overflow-y: auto;
    background: var(--panel);
    border-right: 1px solid var(--border);
    padding: 10px;
  }}
  .thumb {{
    border: 2px solid transparent;
    border-radius: 8px;
    padding: 6px;
    margin-bottom: 10px;
    cursor: pointer;
    transition: border-color .15s ease, background .15s ease;
  }}
  .thumb:hover {{ background: var(--panel-2); }}
  .thumb.active {{ border-color: var(--accent); background: var(--panel-2); }}
  .thumb img {{
    width: 100%;
    display: block;
    border-radius: 4px;
    background: #fff;
  }}
  .thumb .num {{
    font-size: 11px;
    color: var(--text-dim);
    margin-top: 4px;
    text-align: center;
  }}
  #stage-wrap {{
    flex: 1;
    display: flex;
    flex-direction: column;
    min-width: 0;
  }}
  #stage {{
    flex: 1;
    overflow: auto;
    display: flex;
    align-items: center;
    justify-content: center;
    padding: 20px;
    position: relative;
  }}
  #slide-frame {{
    position: relative;
    box-shadow: 0 8px 30px rgba(0,0,0,0.5);
    background: #fff;
    line-height: 0;
  }}
  #slide-frame img {{
    display: block;
    max-width: 100%;
    height: auto;
  }}
  .overlay {{
    position: absolute;
    border: 1.5px dashed rgba(91, 140, 255, 0.55);
    background: rgba(91, 140, 255, 0.06);
    cursor: pointer;
    transition: background .1s ease, border-color .1s ease;
  }}
  .overlay:hover {{
    background: rgba(91, 140, 255, 0.22);
    border-color: var(--accent);
  }}
  .overlay.selected {{
    background: rgba(255, 180, 84, 0.25);
    border: 2px solid var(--accent-2);
  }}
  #info-panel {{
    flex-shrink: 0;
    height: 190px;
    background: var(--panel);
    border-top: 1px solid var(--border);
    padding: 14px 20px;
    overflow-y: auto;
    font-size: 13px;
  }}
  #info-panel h2 {{
    margin: 0 0 8px 0;
    font-size: 13px;
    color: var(--text-dim);
    font-weight: 600;
    text-transform: uppercase;
    letter-spacing: 0.04em;
  }}
  #info-content {{
    color: var(--text-dim);
  }}
  .info-grid {{
    display: grid;
    grid-template-columns: 140px 1fr;
    row-gap: 6px;
    column-gap: 12px;
  }}
  .info-grid .k {{
    color: var(--text-dim);
  }}
  .info-grid .v {{
    color: var(--text);
    font-family: "SF Mono", Menlo, Consolas, monospace;
    font-size: 12.5px;
    word-break: break-word;
    white-space: pre-wrap;
  }}
  .empty-hint {{
    color: var(--text-dim);
    font-style: italic;
  }}
  .badge {{
    display: inline-block;
    padding: 1px 8px;
    border-radius: 999px;
    background: var(--panel-2);
    border: 1px solid var(--border);
    color: var(--accent);
    font-size: 11px;
    margin-left: 8px;
  }}
</style>
</head>
<body>
<header>
  <h1>{title}</h1>
  <span class="meta">{n_slides} slides &middot; mu-ippt interactive_preview v{version} &middot; render: {render_mode}</span>
  <div class="nav">
    <button id="prev-btn">&larr; Prev</button>
    <span id="page-indicator">1 / {n_slides}</span>
    <button id="next-btn">Next &rarr;</button>
  </div>
</header>
<main>
  <div id="sidebar"></div>
  <div id="stage-wrap">
    <div id="stage">
      <div id="slide-frame"></div>
    </div>
    <div id="info-panel">
      <h2>Selected Shape <span id="selected-badge"></span></h2>
      <div id="info-content"><span class="empty-hint">Click any shape outline on the slide to inspect it.</span></div>
    </div>
  </div>
</main>
<script>
const SLIDES = {slides_json};

let currentSlide = 0;
let selectedPath = null;

function fmtEmu(v) {{
  return v.toLocaleString() + " EMU";
}}

function renderSidebar() {{
  const sidebar = document.getElementById('sidebar');
  sidebar.innerHTML = '';
  SLIDES.forEach((slide, i) => {{
    const div = document.createElement('div');
    div.className = 'thumb' + (i === currentSlide ? ' active' : '');
    div.dataset.index = i;
    div.innerHTML = `<img src="${{slide.image}}" alt="Slide ${{i+1}}"><div class="num">${{i+1}}</div>`;
    div.addEventListener('click', () => selectSlide(i));
    sidebar.appendChild(div);
  }});
}}

function selectSlide(i) {{
  currentSlide = i;
  selectedPath = null;
  renderStage();
  renderSidebar();
  renderInfo(null);
  document.getElementById('page-indicator').textContent = `${{i+1}} / ${{SLIDES.length}}`;
  document.getElementById('prev-btn').disabled = (i === 0);
  document.getElementById('next-btn').disabled = (i === SLIDES.length - 1);
}}

function renderStage() {{
  const slide = SLIDES[currentSlide];
  const frame = document.getElementById('slide-frame');
  frame.innerHTML = '';
  frame.style.width = 'auto';

  const img = document.createElement('img');
  img.src = slide.image;
  img.alt = `Slide ${{currentSlide + 1}}`;
  frame.appendChild(img);

  img.addEventListener('load', () => placeOverlays(frame, img, slide), {{ once: true }});
  if (img.complete) placeOverlays(frame, img, slide);
}}

function placeOverlays(frame, img, slide) {{
  const w = img.clientWidth || img.naturalWidth;
  const h = img.clientHeight || img.naturalHeight;
  slide.shapes.forEach((shape) => {{
    if (shape.depth > 0) return; // keep overlay layer to top-level shapes for clarity
    const left = shape.left_emu / slide.slide_width_emu * w;
    const top = shape.top_emu / slide.slide_height_emu * h;
    const width = shape.width_emu / slide.slide_width_emu * w;
    const height = shape.height_emu / slide.slide_height_emu * h;

    const overlay = document.createElement('div');
    overlay.className = 'overlay' + (shape.path === selectedPath ? ' selected' : '');
    overlay.style.left = left + 'px';
    overlay.style.top = top + 'px';
    overlay.style.width = Math.max(width, 2) + 'px';
    overlay.style.height = Math.max(height, 2) + 'px';
    overlay.title = shape.name || shape.shape_type;
    overlay.dataset.path = shape.path;
    overlay.addEventListener('click', (e) => {{
      e.stopPropagation();
      selectedPath = shape.path;
      renderInfo(shape);
      frame.querySelectorAll('.overlay').forEach(o => o.classList.remove('selected'));
      overlay.classList.add('selected');
    }});
    frame.appendChild(overlay);
  }});
}}

function renderInfo(shape) {{
  const content = document.getElementById('info-content');
  const badge = document.getElementById('selected-badge');
  if (!shape) {{
    content.innerHTML = '<span class="empty-hint">Click any shape outline on the slide to inspect it.</span>';
    badge.textContent = '';
    return;
  }}
  badge.textContent = shape.shape_type;
  const rows = [
    ['Path', shape.path],
    ['Name', shape.name || '(unnamed)'],
    ['Type', shape.shape_type],
    ['Placeholder', shape.is_placeholder ? 'yes' : 'no'],
    ['Text', shape.text || '(no text)'],
    ['Position', `left=${{shape.left_in}}in top=${{shape.top_in}}in`],
    ['Size', `width=${{shape.width_in}}in height=${{shape.height_in}}in`],
    ['Position (EMU)', `left=${{fmtEmu(shape.left_emu)}} top=${{fmtEmu(shape.top_emu)}}`],
    ['Size (EMU)', `width=${{fmtEmu(shape.width_emu)}} height=${{fmtEmu(shape.height_emu)}}`],
  ];
  content.innerHTML = '<div class="info-grid">' + rows.map(
    ([k, v]) => `<div class="k">${{k}}</div><div class="v">${{escapeHtml(String(v))}}</div>`
  ).join('') + '</div>';
}}

function escapeHtml(str) {{
  const div = document.createElement('div');
  div.textContent = str;
  return div.innerHTML;
}}

document.getElementById('prev-btn').addEventListener('click', () => {{
  if (currentSlide > 0) selectSlide(currentSlide - 1);
}});
document.getElementById('next-btn').addEventListener('click', () => {{
  if (currentSlide < SLIDES.length - 1) selectSlide(currentSlide + 1);
}});
document.addEventListener('keydown', (e) => {{
  if (e.key === 'ArrowLeft') document.getElementById('prev-btn').click();
  if (e.key === 'ArrowRight') document.getElementById('next-btn').click();
}});
window.addEventListener('resize', renderStage);

renderSidebar();
selectSlide(0);
</script>
</body>
</html>
"""


def build_html(
    pptx_path: Path,
    slide_w_emu: int,
    slide_h_emu: int,
    slides_shapes: List[List[Dict[str, Any]]],
    images: List[Path],
    render_mode: str,
    embed_images: bool = True,
) -> str:
    slides_payload = []
    for idx, shapes in enumerate(slides_shapes):
        image_ref = ""
        if idx < len(images):
            img_path = images[idx]
            image_ref = image_to_data_uri(img_path) if embed_images else img_path.name
        slides_payload.append({
            "index": idx,
            "image": image_ref,
            "slide_width_emu": slide_w_emu,
            "slide_height_emu": slide_h_emu,
            "shapes": shapes,
        })

    slides_json = json.dumps(slides_payload, ensure_ascii=False)
    return HTML_TEMPLATE.format(
        title=html.escape(pptx_path.name),
        n_slides=len(slides_shapes),
        version=__version__,
        render_mode=render_mode,
        slides_json=slides_json,
    )


# --------------------------------------------------------------------------
# Orchestration
# --------------------------------------------------------------------------

def generate_preview(pptx_path: Path, output_html: Optional[Path], dpi: int = 150) -> Path:
    if not pptx_path.exists():
        print(f"Error: {pptx_path} not found", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(pptx_path))
    slide_w, slide_h, slides_shapes = build_slide_index(prs)
    n_slides = len(slides_shapes)

    if output_html is None:
        output_html = pptx_path.parent / f"{pptx_path.stem}_preview.html"
    output_html = Path(output_html)
    output_html.parent.mkdir(parents=True, exist_ok=True)

    render_mode = "unknown"
    images: List[Path] = []

    # Strategy 1: reuse visual_verify.py output
    existing = find_existing_review_images(pptx_path)
    if existing and len(existing) >= n_slides:
        images = existing[:n_slides]
        render_mode = "visual_verify cache"
        print(f"Using {len(images)} existing slide images from visual_review/")
    else:
        # Strategy 2: LibreOffice rasterization
        work_dir = output_html.parent / f".{pptx_path.stem}_preview_assets"
        rendered = render_via_libreoffice(pptx_path, work_dir, dpi=dpi)
        if rendered and len(rendered) >= n_slides:
            images = rendered[:n_slides]
            render_mode = "LibreOffice + pdftoppm"
            print(f"Rendered {len(images)} slide images via LibreOffice")
        else:
            # Strategy 3: pure python-pptx + Pillow placeholder boxes
            if not _HAS_PIL:
                print(
                    "Error: LibreOffice/pdftoppm unavailable and Pillow is not installed; "
                    "cannot generate any preview images.",
                    file=sys.stderr,
                )
                sys.exit(1)
            images = render_placeholders(slide_w, slide_h, slides_shapes, work_dir)
            render_mode = "placeholder (python-pptx + Pillow)"
            print(
                f"LibreOffice unavailable; drew {len(images)} placeholder box renders "
                "(install LibreOffice for accurate visuals)"
            )

    html_content = build_html(
        pptx_path, slide_w, slide_h, slides_shapes, images, render_mode, embed_images=True
    )
    output_html.write_text(html_content, encoding="utf-8")
    print(f"Preview HTML written to: {output_html}")
    print(f"Render mode: {render_mode}")
    return output_html


def serve_preview(html_path: Path, port: int = 0) -> None:
    import http.server
    import socketserver
    import webbrowser

    directory = str(html_path.parent)

    class Handler(http.server.SimpleHTTPRequestHandler):
        def __init__(self, *args, **kwargs):
            super().__init__(*args, directory=directory, **kwargs)

        def log_message(self, format, *args):  # noqa: A002
            pass

    with socketserver.TCPServer(("127.0.0.1", port), Handler) as httpd:
        actual_port = httpd.server_address[1]
        url = f"http://127.0.0.1:{actual_port}/{html_path.name}"
        print(f"Serving preview at: {url}")
        print("Press Ctrl+C to stop.")
        try:
            webbrowser.open(url)
        except Exception:
            pass
        try:
            httpd.serve_forever()
        except KeyboardInterrupt:
            print("\nStopped.")


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Generate an interactive HTML preview for a PPTX with clickable shape overlays.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python3 interactive_preview.py deck.pptx
    Writes deck_preview.html next to the input file.

  python3 interactive_preview.py deck.pptx --output /tmp/preview.html --port 0
    Writes the HTML file and serves it on an OS-assigned local port,
    opening the default browser automatically.
        """,
    )
    parser.add_argument("pptx", help="Input PowerPoint file (.pptx)")
    parser.add_argument("--output", "-o", help="Output HTML file path")
    parser.add_argument("--port", type=int, default=None, help="Serve the preview over HTTP on this port (0 = auto)")
    parser.add_argument("--dpi", type=int, default=150, help="Rasterization DPI for LibreOffice fallback (default: 150)")
    parser.add_argument("--version", action="version", version=f"%(prog)s {__version__}")
    args = parser.parse_args()

    pptx_path = Path(args.pptx).resolve()
    output_html = Path(args.output).resolve() if args.output else None

    html_path = generate_preview(pptx_path, output_html, dpi=args.dpi)

    if args.port is not None:
        serve_preview(html_path, port=args.port)


if __name__ == "__main__":
    main()
