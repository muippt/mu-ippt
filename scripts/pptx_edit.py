#!/usr/bin/env python3
"""
pptx_edit.py — Lightweight post-edit CLI for PPTX (P2 tool)

Purpose:
    Apply small, targeted edits to an already-generated PPTX without
    re-running the full SVG -> DrawingML -> PPTX pipeline. Wraps and
    extends the existing pptx_editing/color_unify.py and
    pptx_editing/replace.py capabilities behind a single subcommand CLI,
    and reuses pptx_inspect.py's addressing scheme and issue detector.

Subcommands:
    set           <pptx> <path> --prop key=value [key=value ...]
    find-replace  <pptx> --find TEXT --replace TEXT [--slide N] [--regex]
    recolor       <pptx> --bg HEX --primary HEX --accent HEX [--dry-run]
    add-slide     <pptx> --after N
    remove-slide  <pptx> N
    move-slide    <pptx> FROM TO

All operations write to a new file (--output); the input file is never
modified in place. After a successful edit, pptx_inspect.py's issue
detector is run automatically on the result unless --skip-check is given.

Usage:
    python3 pptx_edit.py set deck.pptx '/slide[1]/shape[2]' --prop text="New Title" --output out.pptx
    python3 pptx_edit.py find-replace deck.pptx --find "2024" --replace "2025" --output out.pptx
    python3 pptx_edit.py recolor deck.pptx --bg 0D0D2B --primary 00FF88 --accent FF6B35 --output out.pptx
    python3 pptx_edit.py add-slide deck.pptx --after 3 --output out.pptx
    python3 pptx_edit.py remove-slide deck.pptx 5 --output out.pptx
    python3 pptx_edit.py move-slide deck.pptx 5 2 --output out.pptx

Exit codes:
    0 success
    1 error
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

__version__ = "1.2.0"

SCRIPT_DIR = Path(__file__).resolve().parent
PPTX_EDITING_DIR = SCRIPT_DIR / "pptx_editing"

# Make sibling modules importable regardless of current working directory.
for _p in (str(SCRIPT_DIR), str(PPTX_EDITING_DIR)):
    if _p not in sys.path:
        sys.path.insert(0, _p)

# pptx_inspect.py lives alongside this script.
try:
    import pptx_inspect  # noqa: E402
except ImportError:
    pptx_inspect = None  # type: ignore

# color_unify.py lives in pptx_editing/.
try:
    import color_unify  # noqa: E402
except ImportError:
    color_unify = None  # type: ignore


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def load_presentation(pptx_path: Path) -> Any:
    if not pptx_path.exists():
        print(f"Error: file not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)
    try:
        return Presentation(str(pptx_path))
    except Exception as e:
        print(f"Error: could not open PPTX: {e}", file=sys.stderr)
        sys.exit(1)


def default_output_path(input_path: Path, suffix: str) -> Path:
    return input_path.with_name(f"{input_path.stem}.{suffix}{input_path.suffix}")


def save_and_check(prs: Any, output_path: Path, skip_check: bool) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Saved: {output_path}")

    if skip_check or pptx_inspect is None:
        return

    try:
        checked_prs = Presentation(str(output_path))
        issues = pptx_inspect.detect_issues(checked_prs)
        if issues:
            print(f"\n{len(issues)} issue(s) detected after edit:", file=sys.stderr)
            for issue in issues:
                print(f"  {issue['path']} [{issue['issue']}] {issue['detail']}", file=sys.stderr)
        else:
            print("Post-edit check: no issues detected")
    except Exception as e:
        print(f"Warning: post-edit issue check failed: {e}", file=sys.stderr)


def hex_to_rgbcolor(hex_str: str) -> RGBColor:
    h = hex_str.strip().lstrip("#")
    if len(h) != 6:
        raise ValueError(f"Invalid hex color: '{hex_str}'")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ---------------------------------------------------------------------------
# Subcommand: set
# ---------------------------------------------------------------------------

SUPPORTED_SET_PROPS = {"text", "font", "size", "color", "fill", "bold", "italic", "alignment"}

ALIGNMENT_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def parse_prop_args(prop_args: List[str]) -> Dict[str, str]:
    """Parse a list of key=value strings into a dict."""
    props: Dict[str, str] = {}
    for item in prop_args:
        if "=" not in item:
            raise ValueError(f"Invalid --prop value: '{item}'. Expected key=value.")
        key, value = item.split("=", 1)
        key = key.strip().lower()
        if key not in SUPPORTED_SET_PROPS:
            raise ValueError(
                f"Unsupported property: '{key}'. Supported: {', '.join(sorted(SUPPORTED_SET_PROPS))}"
            )
        props[key] = value
    return props


def apply_set_props(shape: Any, props: Dict[str, str]) -> None:
    """Apply a dict of properties to a shape (text/font/size/color/fill/bold/italic/alignment)."""
    has_text_frame = hasattr(shape, "has_text_frame") and shape.has_text_frame

    if "fill" in props:
        color = hex_to_rgbcolor(props["fill"])
        shape.fill.solid()
        shape.fill.fore_color.rgb = color

    text_props = {k: v for k, v in props.items() if k != "fill"}
    if not text_props:
        return

    if not has_text_frame:
        raise ValueError("Shape has no text frame; cannot apply text/font properties")

    text_frame = shape.text_frame

    if "text" in text_props:
        # Replace all text with a single paragraph/run, preserving basic formatting intent.
        new_text = text_props["text"]
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        run = paragraph.add_run() if not paragraph.runs else paragraph.runs[0]
        run.text = new_text

    # Apply remaining formatting to every run in every paragraph.
    for paragraph in text_frame.paragraphs:
        if "alignment" in text_props:
            align_key = text_props["alignment"].strip().lower()
            if align_key not in ALIGNMENT_MAP:
                raise ValueError(
                    f"Invalid alignment: '{text_props['alignment']}'. "
                    f"Supported: {', '.join(ALIGNMENT_MAP)}"
                )
            paragraph.alignment = ALIGNMENT_MAP[align_key]

        if not paragraph.runs:
            continue

        for run in paragraph.runs:
            if "font" in text_props:
                run.font.name = text_props["font"]
            if "size" in text_props:
                try:
                    run.font.size = Pt(float(text_props["size"]))
                except ValueError:
                    raise ValueError(f"Invalid size: '{text_props['size']}'")
            if "bold" in text_props:
                run.font.bold = text_props["bold"].strip().lower() in ("1", "true", "yes")
            if "italic" in text_props:
                run.font.italic = text_props["italic"].strip().lower() in ("1", "true", "yes")
            if "color" in text_props:
                run.font.color.rgb = hex_to_rgbcolor(text_props["color"])


def cmd_set(args: argparse.Namespace) -> int:
    pptx_path = Path(args.pptx)
    prs = load_presentation(pptx_path)

    try:
        props = parse_prop_args(args.prop)
        if pptx_inspect is None:
            print("Error: pptx_inspect.py could not be imported for path resolution", file=sys.stderr)
            return 1
        _slide, shape, _slide_idx, shape_idx = pptx_inspect.resolve_path(prs, args.path)
        if shape is None:
            print(f"Error: path '{args.path}' does not resolve to a shape", file=sys.stderr)
            return 1
        apply_set_props(shape, props)
    except (ValueError, Exception) as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1

    output_path = Path(args.output) if args.output else default_output_path(pptx_path, "edited")
    save_and_check(prs, output_path, args.skip_check)
    return 0


# ---------------------------------------------------------------------------
# Subcommand: find-replace
# ---------------------------------------------------------------------------

def cmd_find_replace(args: argparse.Namespace) -> int:
    pptx_path = Path(args.pptx)
    prs = load_presentation(pptx_path)

    pattern: Optional[re.Pattern] = None
    if args.regex:
        try:
            pattern = re.compile(args.find)
        except re.error as e:
            print(f"Error: invalid regex '{args.find}': {e}", file=sys.stderr)
            return 1

    slides = list(prs.slides)
    target_indices = range(len(slides))
    if args.slide is not None:
        if args.slide < 1 or args.slide > len(slides):
            print(f"Error: --slide {args.slide} out of range (1-{len(slides)})", file=sys.stderr)
            return 1
        target_indices = range(args.slide - 1, args.slide)

    replacements_made = 0

    def replace_in_text(text: str) -> Optional[str]:
        nonlocal replacements_made
        if pattern is not None:
            new_text, count = pattern.subn(args.replace, text)
            if count:
                replacements_made += count
                return new_text
            return None
        if args.find in text:
            count = text.count(args.find)
            replacements_made += count
            return text.replace(args.find, args.replace)
        return None

    def walk_shapes(shapes):
        for shape in shapes:
            if hasattr(shape, "shapes"):  # group shape
                walk_shapes(shape.shapes)
                continue
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        new_text = replace_in_text(run.text)
                        if new_text is not None:
                            run.text = new_text
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                new_text = replace_in_text(run.text)
                                if new_text is not None:
                                    run.text = new_text

    for idx in target_indices:
        walk_shapes(slides[idx].shapes)

    if replacements_made == 0:
        print(f"Warning: no matches found for '{args.find}'", file=sys.stderr)

    output_path = Path(args.output) if args.output else default_output_path(pptx_path, "edited")
    save_and_check(prs, output_path, args.skip_check)
    print(f"Replacements made: {replacements_made}")
    return 0


# ---------------------------------------------------------------------------
# Subcommand: recolor (wraps color_unify.py)
# ---------------------------------------------------------------------------

def cmd_recolor(args: argparse.Namespace) -> int:
    pptx_path = Path(args.pptx)

    if color_unify is None:
        print("Error: could not import pptx_editing/color_unify.py", file=sys.stderr)
        return 1

    theme_map = color_unify.build_theme_map(
        primary=args.primary.upper(),
        accent=args.accent.upper(),
        bg=args.bg.upper(),
        text=args.text.upper(),
        card_bg=args.card_bg.upper() if args.card_bg else None,
    )

    prs = load_presentation(pptx_path)

    if args.dry_run:
        print("Dry run — colors that would be replaced:")
        found: Dict[str, int] = {}
        for slide in prs.slides:
            for shape in slide.shapes:
                for child in shape._element.iter():
                    if child.tag.endswith("}srgbClr"):
                        val = child.get("val", "").upper()
                        if val in theme_map:
                            found[val] = found.get(val, 0) + 1
        if not found:
            print("  (no matching default colors found)")
        for old_val, count in sorted(found.items()):
            new_val = theme_map[old_val]
            print(f"  #{old_val} -> #{new_val}  ({count} occurrence(s))")
        return 0

    bg_color = color_unify.hex_to_rgb(args.bg)
    for slide in prs.slides:
        if args.dark_bg:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = bg_color
        for shape in slide.shapes:
            color_unify.replace_colors(shape._element, theme_map)

    output_path = Path(args.output) if args.output else default_output_path(pptx_path, "recolored")
    save_and_check(prs, output_path, args.skip_check)
    print(f"bg={args.bg} primary={args.primary} accent={args.accent}")
    return 0


# ---------------------------------------------------------------------------
# Subcommand: add-slide
# ---------------------------------------------------------------------------

def cmd_add_slide(args: argparse.Namespace) -> int:
    pptx_path = Path(args.pptx)
    prs = load_presentation(pptx_path)

    slides = prs.slides
    total = len(slides)
    if args.after < 0 or args.after > total:
        print(f"Error: --after {args.after} out of range (0-{total})", file=sys.stderr)
        return 1

    # Use a blank layout if available, otherwise fall back to the layout of
    # the slide right before the insertion point (or the first layout).
    blank_layout = None
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name.strip().lower() == "blank":
            blank_layout = layout
            break
    if blank_layout is None:
        if args.after > 0:
            blank_layout = slides[args.after - 1].slide_layout
        else:
            blank_layout = prs.slide_masters[0].slide_layouts[0]

    new_slide = slides.add_slide(blank_layout)

    # Remove any placeholder shapes copied in from the layout to keep it truly blank.
    for shape in list(new_slide.shapes):
        shape.element.getparent().remove(shape.element)

    # Reposition the newly appended slide to just after `args.after`.
    xml_slides = slides._sldIdLst
    new_slide_elem = xml_slides[-1]
    xml_slides.remove(new_slide_elem)
    xml_slides.insert(args.after, new_slide_elem)

    output_path = Path(args.output) if args.output else default_output_path(pptx_path, "edited")
    save_and_check(prs, output_path, args.skip_check)
    print(f"Inserted blank slide after position {args.after} (new slide is now #{args.after + 1})")
    return 0


# ---------------------------------------------------------------------------
# Subcommand: remove-slide
# ---------------------------------------------------------------------------

def cmd_remove_slide(args: argparse.Namespace) -> int:
    pptx_path = Path(args.pptx)
    prs = load_presentation(pptx_path)

    total = len(prs.slides)
    if args.index < 1 or args.index > total:
        print(f"Error: slide index {args.index} out of range (1-{total})", file=sys.stderr)
        return 1

    idx0 = args.index - 1
    xml_slides = prs.slides._sldIdLst
    rId = xml_slides[idx0].rId
    prs.part.drop_rel(rId)
    del xml_slides[idx0]

    output_path = Path(args.output) if args.output else default_output_path(pptx_path, "edited")
    save_and_check(prs, output_path, args.skip_check)
    print(f"Removed slide {args.index}. Presentation now has {len(prs.slides)} slide(s)")
    return 0


# ---------------------------------------------------------------------------
# Subcommand: move-slide
# ---------------------------------------------------------------------------

def cmd_move_slide(args: argparse.Namespace) -> int:
    pptx_path = Path(args.pptx)
    prs = load_presentation(pptx_path)

    total = len(prs.slides)
    if args.from_index < 1 or args.from_index > total:
        print(f"Error: from-index {args.from_index} out of range (1-{total})", file=sys.stderr)
        return 1
    if args.to_index < 1 or args.to_index > total:
        print(f"Error: to-index {args.to_index} out of range (1-{total})", file=sys.stderr)
        return 1

    xml_slides = prs.slides._sldIdLst
    from0 = args.from_index - 1
    to0 = args.to_index - 1

    slide_elem = xml_slides[from0]
    xml_slides.remove(slide_elem)
    xml_slides.insert(to0, slide_elem)

    output_path = Path(args.output) if args.output else default_output_path(pptx_path, "edited")
    save_and_check(prs, output_path, args.skip_check)
    print(f"Moved slide {args.from_index} -> position {args.to_index}")
    return 0


# ---------------------------------------------------------------------------
# CLI wiring
# ---------------------------------------------------------------------------

def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Lightweight post-edit CLI for PPTX files (no full SVG->PPTX rebuild needed).",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python3 pptx_edit.py set deck.pptx '/slide[1]/shape[2]' --prop text="New Title" --prop bold=true

  python3 pptx_edit.py find-replace deck.pptx --find "2024" --replace "2025" --slide 3

  python3 pptx_edit.py recolor deck.pptx --bg 0D0D2B --primary 00FF88 --accent FF6B35 --dry-run

  python3 pptx_edit.py add-slide deck.pptx --after 3

  python3 pptx_edit.py remove-slide deck.pptx 5

  python3 pptx_edit.py move-slide deck.pptx 5 2

All subcommands accept --output <path> (defaults to <input>.edited.pptx) and
--skip-check to bypass the automatic post-edit issue scan.
        """,
    )
    subparsers = parser.add_subparsers(dest="command", required=True)

    # set
    p_set = subparsers.add_parser("set", help="Modify a single element's properties")
    p_set.add_argument("pptx", help="Input PPTX file")
    p_set.add_argument("path", help="Element path, e.g. /slide[1]/shape[2]")
    p_set.add_argument(
        "--prop", action="append", required=True,
        help="key=value property to set. Repeatable. Keys: text,font,size,color,fill,bold,italic,alignment",
    )
    p_set.add_argument("--output", "-o", help="Output PPTX path (default: <input>.edited.pptx)")
    p_set.add_argument("--skip-check", action="store_true", help="Skip automatic post-edit issue check")
    p_set.set_defaults(func=cmd_set)

    # find-replace
    p_fr = subparsers.add_parser("find-replace", help="Global find/replace across text runs")
    p_fr.add_argument("pptx", help="Input PPTX file")
    p_fr.add_argument("--find", required=True, help="Text (or regex pattern) to find")
    p_fr.add_argument("--replace", required=True, help="Replacement text")
    p_fr.add_argument("--slide", type=int, help="Limit to a single 1-based slide index")
    p_fr.add_argument("--regex", action="store_true", help="Treat --find as a regular expression")
    p_fr.add_argument("--output", "-o", help="Output PPTX path (default: <input>.edited.pptx)")
    p_fr.add_argument("--skip-check", action="store_true", help="Skip automatic post-edit issue check")
    p_fr.set_defaults(func=cmd_find_replace)

    # recolor
    p_rc = subparsers.add_parser("recolor", help="Unify color scheme (wraps color_unify.py)")
    p_rc.add_argument("pptx", help="Input PPTX file")
    p_rc.add_argument("--bg", required=True, help="Background color hex, e.g. 0D0D2B")
    p_rc.add_argument("--primary", required=True, help="Primary color hex, e.g. 00FF88")
    p_rc.add_argument("--accent", required=True, help="Accent color hex, e.g. FF6B35")
    p_rc.add_argument("--text", default="E0E0FF", help="Body text color hex (default: E0E0FF)")
    p_rc.add_argument("--card-bg", default=None, help="Card background hex (default: auto-derived from bg)")
    p_rc.add_argument(
        "--dark-bg", action="store_true", default=True,
        help="Also set every slide's background fill to --bg (default: on)",
    )
    p_rc.add_argument("--dry-run", action="store_true", help="Preview colors that would change, without writing output")
    p_rc.add_argument("--output", "-o", help="Output PPTX path (default: <input>.recolored.pptx)")
    p_rc.add_argument("--skip-check", action="store_true", help="Skip automatic post-edit issue check")
    p_rc.set_defaults(func=cmd_recolor)

    # add-slide
    p_add = subparsers.add_parser("add-slide", help="Insert a blank slide")
    p_add.add_argument("pptx", help="Input PPTX file")
    p_add.add_argument("--after", type=int, required=True, help="1-based slide index to insert after (0 = insert at start)")
    p_add.add_argument("--output", "-o", help="Output PPTX path (default: <input>.edited.pptx)")
    p_add.add_argument("--skip-check", action="store_true", help="Skip automatic post-edit issue check")
    p_add.set_defaults(func=cmd_add_slide)

    # remove-slide
    p_rm = subparsers.add_parser("remove-slide", help="Delete a slide")
    p_rm.add_argument("pptx", help="Input PPTX file")
    p_rm.add_argument("index", type=int, help="1-based slide index to remove")
    p_rm.add_argument("--output", "-o", help="Output PPTX path (default: <input>.edited.pptx)")
    p_rm.add_argument("--skip-check", action="store_true", help="Skip automatic post-edit issue check")
    p_rm.set_defaults(func=cmd_remove_slide)

    # move-slide
    p_mv = subparsers.add_parser("move-slide", help="Move a slide to a new position")
    p_mv.add_argument("pptx", help="Input PPTX file")
    p_mv.add_argument("from_index", type=int, help="1-based current slide index")
    p_mv.add_argument("to_index", type=int, help="1-based target slide index")
    p_mv.add_argument("--output", "-o", help="Output PPTX path (default: <input>.edited.pptx)")
    p_mv.add_argument("--skip-check", action="store_true", help="Skip automatic post-edit issue check")
    p_mv.set_defaults(func=cmd_move_slide)

    return parser


def main() -> int:
    parser = build_parser()
    args = parser.parse_args()
    return args.func(args)


if __name__ == "__main__":
    sys.exit(main())
