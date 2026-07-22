#!/usr/bin/env python3
"""
pptx_inspect.py — Element-level query / addressing system for PPTX (P1 tool)

Purpose:
    Inspect a PPTX file at the element (shape) level, similar to OfficeCLI's
    get/query commands. Built on python-pptx (project's existing dependency).
    Designed to be composed with pptx_edit.py, which uses the same path
    addressing scheme to target elements for modification.

Addressing scheme:
    /slide[N]                 - N is 1-based slide index
    /slide[N]/shape[M]        - M is 1-based shape index within the slide
    /slide[N]/shape[@name=Foo]   - address a shape by its exact name
    /slide[N]/shape[@id=123]     - address a shape by its shape id

Query selectors (used with --query):
    shape[attr=value]     - exact match
    shape[attr!=value]    - not equal
    shape[attr~=text]     - text contains substring (case-insensitive)
    Supported attrs: type, name, fill, text, font, color, bold, italic

Usage:
    python3 pptx_inspect.py <pptx> '/slide[1]'
    python3 pptx_inspect.py <pptx> '/slide[1]/shape[2]'
    python3 pptx_inspect.py <pptx> '/slide[1]/shape[@name=Title 1]'
    python3 pptx_inspect.py <pptx> --json '/slide[1]'
    python3 pptx_inspect.py <pptx> --query 'shape[fill=FF0000]'
    python3 pptx_inspect.py <pptx> outline
    python3 pptx_inspect.py <pptx> stats
    python3 pptx_inspect.py <pptx> issues

Text output format (grep-friendly):
    path (type) "text" key=val key=val ...

Exit codes:
    0 success
    1 error
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

__version__ = "1.2.0"

EMU_PER_INCH = 914400.0


# ---------------------------------------------------------------------------
# Data extraction helpers
# ---------------------------------------------------------------------------

def emu_to_inches(emu: Optional[int]) -> Optional[float]:
    if emu is None:
        return None
    return round(emu / EMU_PER_INCH, 3)


def get_shape_fill_hex(shape: Any) -> Optional[str]:
    """Best-effort extraction of a shape's solid fill color as hex (no '#')."""
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        if str(fill.type).startswith("MSO_FILL_TYPE.SOLID") or fill.type == 1:
            rgb = fill.fore_color.rgb
            if rgb is not None:
                return str(rgb)
    except Exception:
        pass
    return None


def get_shape_text(shape: Any) -> Optional[str]:
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        text = shape.text_frame.text
        return text if text else None
    return None


def get_first_run_font(shape: Any) -> Dict[str, Any]:
    """Extract font name/size/color/bold/italic from the shape's first non-empty run."""
    info: Dict[str, Any] = {}
    if not (hasattr(shape, "has_text_frame") and shape.has_text_frame):
        return info

    for paragraph in shape.text_frame.paragraphs:
        if not paragraph.runs:
            continue
        run = paragraph.runs[0]
        font = run.font
        if font.name:
            info["font"] = font.name
        if font.size:
            info["size"] = font.size.pt
        if font.bold is not None:
            info["bold"] = font.bold
        if font.italic is not None:
            info["italic"] = font.italic
        try:
            if font.color and font.color.rgb:
                info["color"] = str(font.color.rgb)
        except (AttributeError, TypeError):
            pass
        if info:
            break
    return info


def shape_type_name(shape: Any) -> str:
    try:
        return str(shape.shape_type).split(".")[-1].split(" ")[0]
    except Exception:
        return "UNKNOWN"


def describe_shape(shape: Any, slide_idx: int, shape_idx: int) -> Dict[str, Any]:
    """Build a JSON-serializable description of a single shape."""
    path = f"/slide[{slide_idx}]/shape[{shape_idx}]"
    data: Dict[str, Any] = {
        "path": path,
        "id": getattr(shape, "shape_id", None),
        "name": getattr(shape, "name", None),
        "type": shape_type_name(shape),
        "position": {
            "x": emu_to_inches(getattr(shape, "left", None)),
            "y": emu_to_inches(getattr(shape, "top", None)),
            "w": emu_to_inches(getattr(shape, "width", None)),
            "h": emu_to_inches(getattr(shape, "height", None)),
        },
    }

    text = get_shape_text(shape)
    if text is not None:
        data["text"] = text

    fill_hex = get_shape_fill_hex(shape)
    if fill_hex:
        data["fill"] = fill_hex

    font_info = get_first_run_font(shape)
    if font_info:
        data["font_info"] = font_info

    if getattr(shape, "is_placeholder", False):
        try:
            data["placeholder_type"] = str(shape.placeholder_format.type).split(".")[-1].split(" ")[0]
        except Exception:
            pass

    return data


def iter_slide_shapes(slide: Any):
    """Yield (1-based index, shape) pairs for a slide's top-level shapes."""
    for idx, shape in enumerate(slide.shapes, start=1):
        yield idx, shape


# ---------------------------------------------------------------------------
# Path parsing / addressing
# ---------------------------------------------------------------------------

PATH_RE = re.compile(r"^/slide\[(\d+)\](?:/shape\[(.+?)\])?$")


class PathError(ValueError):
    pass


def parse_path(path: str) -> Tuple[int, Optional[str]]:
    """Parse a path like '/slide[1]/shape[2]' or '/slide[1]/shape[@name=Foo]'.

    Returns (slide_index_1based, shape_selector_or_None).
    shape_selector is one of: a digit string, '@name=Value', or '@id=Value'.
    """
    m = PATH_RE.match(path.strip())
    if not m:
        raise PathError(
            f"Invalid path: '{path}'. Expected format: /slide[N] or /slide[N]/shape[M|@name=X|@id=N]"
        )
    slide_idx = int(m.group(1))
    shape_sel = m.group(2)
    return slide_idx, shape_sel


def resolve_shape(slide: Any, selector: str) -> Tuple[int, Any]:
    """Resolve a shape selector against a slide. Returns (1-based idx, shape)."""
    selector = selector.strip()
    if selector.startswith("@name="):
        name = selector.split("=", 1)[1]
        for idx, shape in iter_slide_shapes(slide):
            if getattr(shape, "name", None) == name:
                return idx, shape
        raise PathError(f"No shape found with name='{name}'")
    if selector.startswith("@id="):
        sid = selector.split("=", 1)[1]
        try:
            sid_int = int(sid)
        except ValueError:
            raise PathError(f"Invalid shape id: '{sid}'")
        for idx, shape in iter_slide_shapes(slide):
            if getattr(shape, "shape_id", None) == sid_int:
                return idx, shape
        raise PathError(f"No shape found with id={sid_int}")
    # Plain numeric index
    try:
        idx = int(selector)
    except ValueError:
        raise PathError(f"Invalid shape selector: '{selector}'")
    shapes = list(iter_slide_shapes(slide))
    for si, shape in shapes:
        if si == idx:
            return si, shape
    raise PathError(f"Shape index {idx} out of range (1-{len(shapes)})")


def resolve_path(prs: Any, path: str) -> Tuple[Any, Optional[Any], int, Optional[int]]:
    """Resolve a full path against a Presentation.

    Returns (slide, shape_or_None, slide_idx_1based, shape_idx_1based_or_None).
    """
    slide_idx, shape_sel = parse_path(path)
    slides = list(prs.slides)
    if slide_idx < 1 or slide_idx > len(slides):
        raise PathError(f"Slide index {slide_idx} out of range (1-{len(slides)})")
    slide = slides[slide_idx - 1]

    if shape_sel is None:
        return slide, None, slide_idx, None

    shape_idx, shape = resolve_shape(slide, shape_sel)
    return slide, shape, slide_idx, shape_idx


# ---------------------------------------------------------------------------
# Text output formatting (grep-friendly)
# ---------------------------------------------------------------------------

def format_shape_line(desc: Dict[str, Any]) -> str:
    """Format a shape description as: path (type) "text" key=val key=val ..."""
    parts = [desc["path"], f"({desc['type']})"]

    text = desc.get("text")
    if text:
        preview = text.replace("\n", "\\n")
        if len(preview) > 60:
            preview = preview[:57] + "..."
        parts.append(f'"{preview}"')
    else:
        parts.append('""')

    pos = desc.get("position", {})
    if all(v is not None for v in pos.values()):
        parts.append(f"x={pos['x']} y={pos['y']} w={pos['w']} h={pos['h']}")

    if desc.get("name"):
        parts.append(f"name={desc['name']}")
    if desc.get("id") is not None:
        parts.append(f"id={desc['id']}")
    if desc.get("fill"):
        parts.append(f"fill={desc['fill']}")
    font_info = desc.get("font_info", {})
    if font_info.get("font"):
        parts.append(f"font={font_info['font']}")
    if font_info.get("size") is not None:
        parts.append(f"size={font_info['size']}")
    if font_info.get("color"):
        parts.append(f"color={font_info['color']}")
    if font_info.get("bold") is not None:
        parts.append(f"bold={font_info['bold']}")
    if font_info.get("italic") is not None:
        parts.append(f"italic={font_info['italic']}")
    if desc.get("placeholder_type"):
        parts.append(f"placeholder={desc['placeholder_type']}")

    return " ".join(parts)


# ---------------------------------------------------------------------------
# Query mode: shape[attr=value] / [attr!=value] / [attr~=text]
# ---------------------------------------------------------------------------

QUERY_RE = re.compile(r"^shape\[(\w+)\s*(=|!=|~=)\s*(.*?)\]$")


def parse_query(query: str) -> Tuple[str, str, str]:
    m = QUERY_RE.match(query.strip())
    if not m:
        raise PathError(
            f"Invalid query: '{query}'. Expected format: shape[attr=value|attr!=value|attr~=text]"
        )
    attr, op, value = m.groups()
    return attr, op, value


def shape_matches(desc: Dict[str, Any], attr: str, op: str, value: str) -> bool:
    if attr == "type":
        actual = desc.get("type", "")
    elif attr == "name":
        actual = desc.get("name") or ""
    elif attr == "fill":
        actual = desc.get("fill") or ""
    elif attr == "text":
        actual = desc.get("text") or ""
    elif attr == "font":
        actual = desc.get("font_info", {}).get("font") or ""
    elif attr == "color":
        actual = desc.get("font_info", {}).get("color") or ""
    elif attr == "bold":
        actual = str(desc.get("font_info", {}).get("bold"))
    elif attr == "italic":
        actual = str(desc.get("font_info", {}).get("italic"))
    else:
        raise PathError(f"Unsupported query attribute: '{attr}'")

    actual_str = str(actual)
    if op == "=":
        return actual_str.lower() == value.lower()
    if op == "!=":
        return actual_str.lower() != value.lower()
    if op == "~=":
        return value.lower() in actual_str.lower()
    return False


def run_query(prs: Any, query: str) -> List[Dict[str, Any]]:
    attr, op, value = parse_query(query)
    results = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in iter_slide_shapes(slide):
            desc = describe_shape(shape, slide_idx, shape_idx)
            if shape_matches(desc, attr, op, value):
                results.append(desc)
    return results


# ---------------------------------------------------------------------------
# Views: outline / stats / issues
# ---------------------------------------------------------------------------

def build_outline(prs: Any) -> List[Dict[str, Any]]:
    outline = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        entry: Dict[str, Any] = {"slide": slide_idx, "shapes": []}
        for shape_idx, shape in iter_slide_shapes(slide):
            desc = describe_shape(shape, slide_idx, shape_idx)
            summary = {
                "path": desc["path"],
                "type": desc["type"],
                "name": desc.get("name"),
            }
            if desc.get("text"):
                text = desc["text"].strip().split("\n")[0]
                summary["text"] = text[:60] + ("..." if len(text) > 60 else "")
            entry["shapes"].append(summary)
        outline.append(entry)
    return outline


def build_stats(prs: Any) -> Dict[str, Any]:
    total_shapes = 0
    total_text_chars = 0
    shape_type_counts: Dict[str, int] = {}
    slides_with_text = 0

    for slide in prs.slides:
        slide_has_text = False
        for _, shape in iter_slide_shapes(slide):
            total_shapes += 1
            t = shape_type_name(shape)
            shape_type_counts[t] = shape_type_counts.get(t, 0) + 1
            text = get_shape_text(shape)
            if text:
                total_text_chars += len(text)
                slide_has_text = True
        if slide_has_text:
            slides_with_text += 1

    return {
        "slide_count": len(prs.slides),
        "shape_count": total_shapes,
        "text_char_count": total_text_chars,
        "slides_with_text": slides_with_text,
        "shape_type_counts": shape_type_counts,
        "slide_width_in": emu_to_inches(prs.slide_width),
        "slide_height_in": emu_to_inches(prs.slide_height),
    }


def detect_issues(prs: Any) -> List[Dict[str, Any]]:
    """Detect common problems: empty text boxes, potential overflow, off-slide shapes."""
    issues: List[Dict[str, Any]] = []
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in iter_slide_shapes(slide):
            path = f"/slide[{slide_idx}]/shape[{shape_idx}]"

            # Empty text frame (placeholder with no content, non-title)
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                text = shape.text_frame.text
                if not text.strip() and not getattr(shape, "is_placeholder", False):
                    issues.append({
                        "path": path,
                        "issue": "empty_text_box",
                        "detail": "Shape has a text frame but no text content",
                    })

            left = getattr(shape, "left", None)
            top = getattr(shape, "top", None)
            width = getattr(shape, "width", None)
            height = getattr(shape, "height", None)

            if None not in (left, top, width, height, slide_w, slide_h):
                # Off-slide / overflow risk
                if left < 0 or top < 0:
                    issues.append({
                        "path": path,
                        "issue": "negative_position",
                        "detail": f"Shape positioned at x={emu_to_inches(left)}in, y={emu_to_inches(top)}in (off-slide)",
                    })
                if left + width > slide_w:
                    overflow = emu_to_inches((left + width) - slide_w)
                    issues.append({
                        "path": path,
                        "issue": "overflow_risk_right",
                        "detail": f"Shape extends {overflow}in beyond right edge",
                    })
                if top + height > slide_h:
                    overflow = emu_to_inches((top + height) - slide_h)
                    issues.append({
                        "path": path,
                        "issue": "overflow_risk_bottom",
                        "detail": f"Shape extends {overflow}in beyond bottom edge",
                    })

            # Tiny font size warning
            font_info = get_first_run_font(shape)
            if font_info.get("size") is not None and font_info["size"] < 8:
                issues.append({
                    "path": path,
                    "issue": "tiny_font",
                    "detail": f"Font size {font_info['size']}pt may be too small to read",
                })

    return issues


# ---------------------------------------------------------------------------
# Output rendering
# ---------------------------------------------------------------------------

def print_shape_descriptions(descs: List[Dict[str, Any]], as_json: bool) -> None:
    if as_json:
        print(json.dumps(descs, indent=2, ensure_ascii=False))
    else:
        for desc in descs:
            print(format_shape_line(desc))


def print_outline(outline: List[Dict[str, Any]], as_json: bool) -> None:
    if as_json:
        print(json.dumps(outline, indent=2, ensure_ascii=False))
        return
    for entry in outline:
        print(f"/slide[{entry['slide']}]")
        for shape in entry["shapes"]:
            label = shape.get("text") or shape.get("name") or ""
            print(f"  {shape['path']} ({shape['type']}) {label}")


def print_stats(stats: Dict[str, Any], as_json: bool) -> None:
    if as_json:
        print(json.dumps(stats, indent=2, ensure_ascii=False))
        return
    print(f"slides={stats['slide_count']} shapes={stats['shape_count']} "
          f"text_chars={stats['text_char_count']} slides_with_text={stats['slides_with_text']}")
    print(f"slide_size={stats['slide_width_in']}in x {stats['slide_height_in']}in")
    print("shape_type_counts:")
    for t, count in sorted(stats["shape_type_counts"].items(), key=lambda kv: -kv[1]):
        print(f"  {t}: {count}")


def print_issues(issues: List[Dict[str, Any]], as_json: bool) -> None:
    if as_json:
        print(json.dumps(issues, indent=2, ensure_ascii=False))
        return
    if not issues:
        print("No issues detected")
        return
    for issue in issues:
        print(f"{issue['path']} [{issue['issue']}] {issue['detail']}")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> int:
    parser = argparse.ArgumentParser(
        description="Element-level query and inspection tool for PPTX files.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python3 pptx_inspect.py deck.pptx '/slide[1]'
    List all shapes on slide 1

  python3 pptx_inspect.py deck.pptx '/slide[1]/shape[2]'
    Show details for shape 2 on slide 1

  python3 pptx_inspect.py deck.pptx '/slide[1]/shape[@name=Title 1]' --json
    Show shape by name, as JSON

  python3 pptx_inspect.py deck.pptx --query 'shape[fill=FF0000]'
    Find all shapes with a red fill anywhere in the deck

  python3 pptx_inspect.py deck.pptx outline
    Print a document outline of all slides and shapes

  python3 pptx_inspect.py deck.pptx stats
    Print summary statistics (slide count, shape count, text volume...)

  python3 pptx_inspect.py deck.pptx issues
    Detect common problems (empty text boxes, overflow risk, tiny fonts)
        """,
    )
    parser.add_argument("pptx", help="Input PPTX file")
    parser.add_argument(
        "target",
        nargs="?",
        help="Path expression (e.g. /slide[1]/shape[2]) or view name (outline|stats|issues)",
    )
    parser.add_argument("--query", help="Query selector, e.g. \"shape[fill=FF0000]\"")
    parser.add_argument("--json", action="store_true", help="Output structured JSON instead of text")

    args = parser.parse_args()

    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        print(f"Error: file not found: {pptx_path}", file=sys.stderr)
        return 1

    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        print(f"Error: could not open PPTX: {e}", file=sys.stderr)
        return 1

    try:
        if args.query:
            results = run_query(prs, args.query)
            print_shape_descriptions(results, args.json)
            return 0

        if args.target is None:
            print("Error: provide a path expression, a view name (outline|stats|issues), or --query", file=sys.stderr)
            return 1

        if args.target == "outline":
            print_outline(build_outline(prs), args.json)
            return 0
        if args.target == "stats":
            print_stats(build_stats(prs), args.json)
            return 0
        if args.target == "issues":
            issues = detect_issues(prs)
            print_issues(issues, args.json)
            return 0

        # Path expression
        slide, shape, slide_idx, shape_idx = resolve_path(prs, args.target)

        if shape is not None:
            desc = describe_shape(shape, slide_idx, shape_idx)
            print_shape_descriptions([desc], args.json)
        else:
            # List all shapes on the slide
            descs = [describe_shape(s, slide_idx, i) for i, s in iter_slide_shapes(slide)]
            print_shape_descriptions(descs, args.json)

        return 0

    except PathError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    sys.exit(main())
